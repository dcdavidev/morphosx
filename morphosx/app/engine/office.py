import io
from typing import Tuple
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image, ImageDraw, ImageFont


class OfficeProcessor:
    """
    Engine for generating previews/thumbnails for Office Documents.
    
    Since pure Python rendering of Office to perfect PDF/Image is extremely 
    complex, we generate a 'Summary Card' image.
    """

    def render_thumbnail(self, doc_data: bytes, filename: str) -> bytes:
        """
        Generate a summary image for a DOCX, PPTX or XLSX file.
        """
        ext = filename.split(".")[-1].lower()
        title = "Office Document"
        summary = ""

        try:
            if ext == "docx":
                doc = Document(io.BytesIO(doc_data))
                title = "Word Document"
                # Extract first 5 paragraphs
                summary = "
".join([p.text for p in doc.paragraphs[:5] if p.text.strip()])
            elif ext == "pptx":
                prs = Presentation(io.BytesIO(doc_data))
                title = f"PowerPoint ({len(prs.slides)} slides)"
                if len(prs.slides) > 0:
                    slide = prs.slides[0]
                    summary = "
".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            elif ext == "xlsx":
                wb = load_workbook(io.BytesIO(doc_data), data_only=True)
                ws = wb.active
                title = f"Excel Spreadsheet ({ws.title})"
                # Extract a 5x5 grid
                rows = []
                for row in ws.iter_rows(max_row=10, max_col=5):
                    rows.append(" | ".join([str(cell.value or "") for cell in row]))
                summary = "
".join(rows)

            return self._create_summary_card(title, summary[:500])
            
        except Exception as e:
            return self._create_summary_card("Error", f"Could not parse office file: {str(e)}")

    def _create_summary_card(self, title: str, text: str) -> bytes:
        """Render a text-based summary card image."""
        width, height = 800, 600
        img = Image.new("RGB", (width, height), color=(240, 240, 240))
        draw = ImageDraw.Draw(img)
        
        # Draw a header bar
        draw.rectangle([0, 0, width, 60], fill=(43, 108, 176))
        
        # Simple text drawing (fallback to default if font missing)
        try:
            # We don't want to rely on specific paths, use default
            draw.text((20, 15), title, fill=(255, 255, 255))
            draw.text((20, 80), text, fill=(50, 50, 50))
        except Exception:
            pass

        output = io.BytesIO()
        img.save(output, format="JPEG")
        return output.getvalue()
